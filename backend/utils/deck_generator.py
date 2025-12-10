from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
from utils.logger import logger

"""
Only basic styling is applied here. Bullet points, headers, and table styles.
Bullets are needs to be only in list form.
"""

def add_styled_table(slide, columns, rows):
    """Insert a formatted table into slide"""
    row_count = len(rows) + 1                          # +1 header row
    col_count = len(columns)

    table_shape = slide.shapes.add_table(row_count, col_count, Inches(1), Inches(2), Inches(8), Inches(1.5))
    table = table_shape.table

    # Header Style
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50, 50, 50)         # darker text
        cell.fill.solid()        
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230)  # light gray professional header

    # Data rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(50, 50, 50)         # darker text

    # Column width (optional tuning)
    for col in table.columns:
        col.width = Inches(4)


# def generate_deck(content, deck_file_path):
#     try:
#         # --- Global Theme Styling ---
#         TITLE_FONT_SIZE = 44
#         SUBTITLE_FONT_SIZE = 26
#         BODY_FONT_SIZE = 24
#         PRIMARY_COLOR = RGBColor(30, 50, 120)      # deep corporate blue
#         ACCENT_COLOR = RGBColor(230, 230, 230)     # light grey
#         TEXT_DARK = RGBColor(40, 40, 40)

#         prs = Presentation()
#         content = '''{
#     "title_slide": {"title":"Expansion Opportunity","subtitle":"Expansion Opportunity"},
#     "slides": [
#         {"title": "Executive Summary", "bullets": ["Overconsumption of Resin-A is causing potential revenue loss, while a competitor dependency on Catalyst-X may be impacted by their own supply constraints."
#             ]
#         },
#         {"title": "Commercial Insight", "bullets": [
#                 "This development could result in lost sales and revenue, especially if we're unable to optimize our internal production to meet demand."
#             ]
#         },
#         {"title": "Detected Patterns", "bullets": [],
#             "visual": {
#                 "type": "table",
#                 "columns": ["Pattern", "Impact"],
#                 "rows": [
#                     ["Overconsumption pattern", "Estimated revenue loss"],
#                     ["Temporary volume spikes (no specific SKU)", "Volume-driven revenue increase"]
#                 ]
#             }
#         },
#         {"title": "Revenue Upside Summary", "bullets": ["A estimated monthly revenue leakage of USD 0, with potential for optimization to drive increased sales."
#             ]
#         }
#     ]
# }
# '''
#         content = json.loads(content)
#         title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#         title_slide.shapes.title.text = content["title_slide"]["title"]
#         title_slide.placeholders[1].text = content["title_slide"]["subtitle"]

#         # Title slide styling
#         title = title_slide.shapes.title.text_frame
#         title.paragraphs[0].font.size = Pt(54)
#         title.paragraphs[0].font.bold = True
#         title.paragraphs[0].font.color.rgb = PRIMARY_COLOR

#         subtitle = title_slide.placeholders[1].text_frame
#         subtitle.paragraphs[0].font.size = Pt(32)
#         subtitle.paragraphs[0].font.color.rgb = TEXT_DARK

#         body = slide.placeholders[1].text_frame
#         body.clear()  # clean default text

#         for slide_data in content["slides"]:
#             slide = prs.slides.add_slide(prs.slide_layouts[1])
#             title_shape = slide.shapes.title
#             title_shape.text = slide_data["title"]

#             title_frame = title_shape.text_frame
#             title_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)
#             title_frame.paragraphs[0].font.bold = True
#             title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

#             body = slide.placeholders[1].text_frame
#             body.clear()  # clean default text

#             for bullet in slide_data["bullets"]:
#                 p = body.add_paragraph()
#                 if isinstance(bullet, dict):
#                     text = " | ".join([f"{k}: {v}" for k, v in bullet.items()])
#                     p.text = text
#                 else:
#                     p.text = bullet

#                 p.level = 0

#             if "visual" in slide_data and slide_data["visual"]["type"] == "table":
#                 columns = slide_data["visual"]["columns"]
#                 rows = slide_data["visual"]["rows"]
#                 add_styled_table(slide, columns, rows)
            
        
#         prs.save(deck_file_path)
#         logger.info(f"Deck generated at: {deck_file_path}")
#         return deck_file_path
        
#     except Exception as e:
#         print(f"Error while generate_deck(): {str(e)}")

def generate_deck(content, deck_file_path):
    try:
        prs = Presentation()

        # --- Global Theme Styling ---
        TITLE_FONT_SIZE = 44
        SUBTITLE_FONT_SIZE = 26
        BODY_FONT_SIZE = 24
        PRIMARY_COLOR = RGBColor(30, 50, 120)      # deep corporate blue
        ACCENT_COLOR = RGBColor(230, 230, 230)     # light grey
        TEXT_DARK = RGBColor(40, 40, 40)

        # ---------- JSON Input -----------
        # content = '''{
        #     "title_slide": {"title":"Expansion Opportunity","subtitle":"SKU Pattern Analysis Report"},
        #     "slides": [
        #         {
        #             "title": "Executive Summary",
        #             "bullets": [
        #                 "Overconsumption of Resin-A is causing potential revenue loss",
        #                 "Dependency risk identified on Catalyst-X supplier"
        #             ],
        #             "visual": {
        #                 "type": "bar_chart",
        #                 "data_source": "resin_a_consumption.csv",
        #                 "x_axis": ["Jan", "Feb", "Mar", "Apr", "May"],
        #                 "y_axis": [120, 150, 180, 220, 300]
        #             }
        #         },
        #         {
        #             "title": "Commercial Insight",
        #             "bullets": [
        #                 "Possible missed revenue opportunity due to inefficient production scaling",
        #                 "Strong demand signals offer upsell potential"
        #             ]
        #         },
        #         {
        #             "title": "Detected Patterns",
        #             "bullets": [],
        #             "visual": {
        #                 "type": "table",
        #                 "columns": ["Pattern", "Impact"],
        #                 "rows": [
        #                     ["Overconsumption pattern", "Estimated revenue loss"],
        #                     ["Volume Spikes", "High Revenue Potential"]
        #                 ]
        #             }
        #         },
        #         {
        #             "title": "Revenue Upside Summary",
        #             "bullets": [
        #                 "Upsell opportunity: USD 50K/month forecast",
        #                 "SKU bundling strategy recommended"
        #             ]
        #         }
        #     ]
        # }'''

        # content = json.loads(content)

        # ========= Title Slide ==========
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = content["title_slide"]["title"]
        title_slide.placeholders[1].text = content["title_slide"]["subtitle"]

        # Title slide styling
        title = title_slide.shapes.title.text_frame
        title.paragraphs[0].font.size = Pt(54)
        title.paragraphs[0].font.bold = True
        title.paragraphs[0].font.color.rgb = PRIMARY_COLOR

        subtitle = title_slide.placeholders[1].text_frame
        subtitle.paragraphs[0].font.size = Pt(32)
        subtitle.paragraphs[0].font.color.rgb = TEXT_DARK

        # ========= Content Slides ==========
        for slide_data in content["slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Title formatting
            title_frame = title_shape.text_frame
            title_frame.paragraphs[0].font.size = Pt(TITLE_FONT_SIZE)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

            body = slide.placeholders[1].text_frame
            body.clear()  # clean default text

            # Bullet styling
            for i, bullet in enumerate(slide_data["bullets"]):
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(BODY_FONT_SIZE)
                p.font.color.rgb = TEXT_DARK

            # Table block
            if "visual" in slide_data and slide_data["visual"]["type"] == "table":
                add_styled_table(
                    slide, 
                    slide_data["visual"]["columns"], 
                    slide_data["visual"]["rows"]
                )

            # Add subtle footer
            footer = slide.shapes.add_textbox(Inches(0.2), Inches(7), Inches(9), Inches(0.5))
            footer_tf = footer.text_frame
            footer_tf.text = "Generated by Agentic AI | Confidential"
            footer_tf.paragraphs[0].font.size = Pt(14)
            footer_tf.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
            footer_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        prs.save(deck_file_path)
        logger.info(f"Deck generated at: {deck_file_path}")        
        return deck_file_path

    except Exception as e:
        print(f"Error while generate_deck(): {str(e)}")

